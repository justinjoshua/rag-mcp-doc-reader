"""Ingest documents: parse -> chunk -> embed -> store.

Usage:
    python -m docreader.ingest              # ingest everything under DOCS_DIR
    python -m docreader.ingest --reset      # wipe the collection first
    python -m docreader.ingest --docs ./my_docs
"""
from __future__ import annotations

import argparse
import hashlib
import sys
from pathlib import Path

# Windows console (cp1252) raises UnicodeEncodeError on non-representable chars;
# force UTF-8 so printing filenames/summaries can't crash ingest.
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except (AttributeError, ValueError):
    pass

from .chunking import chunk_text
from .config import CONFIG
from .embeddings import get_embedder
from .store import get_store

TEXT_EXTENSIONS = {".md", ".txt", ".rst", ".markdown", ".csv", ".json"}
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx"} | TEXT_EXTENSIONS
EMBED_BATCH = 64


def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def _read_docx(path: Path) -> str:
    import docx

    doc = docx.Document(str(path))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"# Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            parts.append("Notes: " + slide.notes_slide.notes_text_frame.text)
        if len(parts) > 1:
            out.append("\n".join(parts))
    return "\n\n".join(out)


def parse_file(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf":
        text = _read_pdf(path)
        # No text layer (scanned/image PDF) -> OCR via Gemini if configured.
        if len(text.strip()) < 100 and CONFIG.ocr and CONFIG.gemini_api_key:
            from .ocr import ocr_pdf

            ocr_text = ocr_pdf(path)
            if len(ocr_text.strip()) > len(text.strip()):
                return ocr_text
        return text
    if ext == ".docx":
        return _read_docx(path)
    if ext == ".pptx":
        return _read_pptx(path)
    if ext in TEXT_EXTENSIONS:
        return path.read_text(encoding="utf-8", errors="replace")
    return ""  # unsupported type -> skipped


def _chunk_id(source: str, index: int, text: str) -> str:
    h = hashlib.sha256(f"{source}\x00{index}\x00{text}".encode("utf-8")).hexdigest()
    return h[:32]


def ingest_file(path: Path, source: str | None = None) -> int:
    """Parse -> chunk -> embed -> store one file. Returns chunks added.

    Reused by the CLI directory walk and by the web upload endpoint.
    """
    store = get_store()
    embedder = get_embedder()
    text = parse_file(path)
    if not text.strip():
        return 0
    source = source or path.name
    chunks = chunk_text(text)
    if not chunks:
        return 0

    ids, docs, metas = [], [], []
    for i, chunk in enumerate(chunks):
        ids.append(_chunk_id(source, i, chunk))
        docs.append(chunk)
        metas.append({"source": source, "chunk_index": i})

    # embed + store in batches to keep memory bounded on large files
    for start in range(0, len(docs), EMBED_BATCH):
        batch = docs[start : start + EMBED_BATCH]
        store.add(
            ids=ids[start : start + EMBED_BATCH],
            embeddings=embedder.embed_documents(batch),
            documents=batch,
            metadatas=metas[start : start + EMBED_BATCH],
        )
    return len(chunks)


def ingest(docs_dir: str, reset: bool = False) -> None:
    store = get_store()
    if reset:
        store.reset()
        print("Collection reset.")

    root = Path(docs_dir).expanduser().resolve()
    if not root.exists():
        raise SystemExit(f"Docs directory not found: {root}")

    files = [p for p in root.rglob("*") if p.is_file()]
    total_chunks, parsed_files = 0, 0

    for path in files:
        source = str(path.relative_to(root))
        n = ingest_file(path, source=source)
        if n == 0:
            if path.suffix.lower() in {".pdf", ".docx"} | TEXT_EXTENSIONS:
                print(f"  ! no extractable text: {source}")
            continue
        parsed_files += 1
        total_chunks += n
        print(f"  + {source}: {n} chunks")

    print(
        f"\nDone. Parsed {parsed_files} file(s), added {total_chunks} chunk(s). "
        f"Collection now holds {store.count()} chunk(s)."
    )


def main() -> None:
    ap = argparse.ArgumentParser(description="Ingest documents into the vector store.")
    ap.add_argument("--docs", default=CONFIG.docs_dir, help="Directory of documents.")
    ap.add_argument("--reset", action="store_true", help="Wipe the collection first.")
    args = ap.parse_args()
    ingest(args.docs, reset=args.reset)


if __name__ == "__main__":
    main()
